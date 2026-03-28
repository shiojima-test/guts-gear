#!/usr/bin/env python3
"""Create PowerPoint presentation for 環境マッピングマップ guide"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

NAVY = RGBColor(0x1E, 0x30, 0x60)
BLUE2 = RGBColor(0x2F, 0x6E, 0xB5)
ACCENT = RGBColor(0x5B, 0x9B, 0xD5)
TEAL = RGBColor(0x3A, 0xAB, 0xA8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x24, 0x33)
GRAY = RGBColor(0x4A, 0x60, 0x80)
LIGHT_GRAY = RGBColor(0x8F, 0xAA, 0xBB)

SS = "/home/user/guts-gear/screenshots"
OUT = "/home/user/guts-gear/mapping_guide.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height


def add_bg(slide, color=NAVY):
    """Fill slide background"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_para(tf, text, size=14, color=WHITE, bold=False, space_before=Pt(4), space_after=Pt(2), align=PP_ALIGN.LEFT, font_name="Meiryo"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    if space_before:
        p.space_before = space_before
    if space_after:
        p.space_after = space_after
    return p


# ═══════════════════════════════════════════
# Slide 1: Title
# ═══════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, NAVY)

# Accent stripe top
add_shape(slide1, 0, 0, SW, Inches(0.08), fill_color=TEAL)

# Accent line left
add_shape(slide1, Inches(1.2), Inches(1.5), Inches(0.06), Inches(3.5), fill_color=ACCENT)

# Title
add_text(slide1, Inches(1.6), Inches(1.8), Inches(10), Inches(1),
         "国立青少年教育施設", size=20, color=ACCENT, bold=False)
add_text(slide1, Inches(1.6), Inches(2.4), Inches(10), Inches(1.2),
         "環境マッピングマップ", size=44, color=WHITE, bold=True)
add_text(slide1, Inches(1.6), Inches(3.4), Inches(10), Inches(0.8),
         "使い方ガイド", size=36, color=ACCENT, bold=False)

# Subtitle
add_text(slide1, Inches(1.6), Inches(4.8), Inches(8), Inches(0.5),
         "高校情報II データサイエンス授業用  |  環境センサーデータの可視化・分析ツール",
         size=14, color=LIGHT_GRAY)

# Bottom bar
add_shape(slide1, 0, SH - Inches(0.6), SW, Inches(0.6), fill_color=BLUE2)
add_text(slide1, Inches(1), SH - Inches(0.55), Inches(11), Inches(0.5),
         "https://shiojima-test.github.io/schoomy-gis/niye_comfort_v1.html",
         size=11, color=WHITE, align=PP_ALIGN.CENTER)

print("Slide 1: Title done")


# ═══════════════════════════════════════════
# Slide 2: Full screen overview with numbered areas
# ═══════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, NAVY)
add_shape(slide2, 0, 0, SW, Inches(0.08), fill_color=TEAL)

add_text(slide2, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
         "画面構成の全体像", size=28, color=WHITE, bold=True)
add_text(slide2, Inches(0.5), Inches(0.65), Inches(12), Inches(0.4),
         "各エリアの役割を理解しましょう", size=14, color=LIGHT_GRAY)

# Full screenshot
img_path = f"{SS}/01_full.png"
slide2.shapes.add_picture(img_path, Inches(0.3), Inches(1.2), Inches(9.2), Inches(5.9))

# Numbered callout boxes on the right
callouts = [
    ("1", "ヘッダー", "タイトル、シリアル接続\nボタン、動作モード表示"),
    ("2", "コントロールバー", "季節（春夏秋冬）と\n時間帯（朝昼夕夜）の切替"),
    ("3", "施設一覧（左パネル）", "29施設のスコア一覧\nクリックで地図移動"),
    ("4", "地図エリア", "マーカーで快適スコアを\n色分け表示"),
    ("5", "データ分析（右パネル）", "季節/分布/相関/時系列\n/異常値/CSV出力"),
    ("6", "凡例", "快適スコアの色分けと\n施設種別の表示"),
]

cx = Inches(9.8)
cy = Inches(1.1)
for num, title, desc in callouts:
    # Number circle
    circle = slide2.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = TEAL
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Meiryo"

    add_text(slide2, cx + Inches(0.5), cy - Inches(0.02), Inches(2.8), Inches(0.3),
             title, size=13, color=ACCENT, bold=True)
    add_text(slide2, cx + Inches(0.5), cy + Inches(0.25), Inches(2.8), Inches(0.6),
             desc, size=10, color=LIGHT_GRAY)
    cy += Inches(0.95)

print("Slide 2: Overview done")


# ═══════════════════════════════════════════
# Slide 3: How to use STEP 1-4
# ═══════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, NAVY)
add_shape(slide3, 0, 0, SW, Inches(0.08), fill_color=TEAL)

add_text(slide3, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
         "使い方ガイド", size=28, color=WHITE, bold=True)
add_text(slide3, Inches(0.5), Inches(0.65), Inches(12), Inches(0.4),
         "4つのステップで環境データを探索しましょう", size=14, color=LIGHT_GRAY)

steps = [
    ("STEP 1", "URLにアクセス", "ブラウザで以下のURLを開きます。\nChrome / Edge 推奨。\nデモモードで自動的にデータが表示されます。",
     "https://shiojima-test.github.io/\nschoomy-gis/niye_comfort_v1.html"),
    ("STEP 2", "季節・時間帯を選択", "コントロールバーで季節（春/夏/秋/冬）と\n時間帯（朝/昼/夕/夜）を切り替えます。\n全16パターンのデータを比較できます。", None),
    ("STEP 3", "施設を探索", "左パネルの施設一覧から施設をクリック、\nまたは地図上のマーカーをクリックします。\nポップアップで詳細データを確認できます。", None),
    ("STEP 4", "データを分析", "右パネルのタブを使って分析します。\n分布（ヒストグラム/箱ひげ図）\n相関（散布図/回帰分析）\n時系列/異常値検知/CSV出力", None),
]

for i, (step_num, title, desc, extra) in enumerate(steps):
    sx = Inches(0.5) + Inches(3.15) * i
    sy = Inches(1.3)

    # Card background
    card = add_rounded_rect(slide3, sx, sy, Inches(3.0), Inches(5.6), fill_color=RGBColor(0x16, 0x24, 0x4A))

    # Step number badge
    badge = add_rounded_rect(slide3, sx + Inches(0.3), sy + Inches(0.3), Inches(1.2), Inches(0.4), fill_color=TEAL)
    tf = badge.text_frame
    tf.paragraphs[0].text = step_num
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Meiryo"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    add_text(slide3, sx + Inches(0.3), sy + Inches(0.9), Inches(2.4), Inches(0.5),
             title, size=18, color=WHITE, bold=True)

    # Description
    add_text(slide3, sx + Inches(0.3), sy + Inches(1.5), Inches(2.4), Inches(2.5),
             desc, size=12, color=LIGHT_GRAY)

    # Add screenshot for specific steps
    if i == 1:
        slide3.shapes.add_picture(f"{SS}/04_header.png", sx + Inches(0.15), sy + Inches(3.2), Inches(2.7), Inches(0.9))
    elif i == 2:
        slide3.shapes.add_picture(f"{SS}/05_popup.png", sx + Inches(0.15), sy + Inches(3.2), Inches(2.7), Inches(2.0))
    elif i == 3:
        slide3.shapes.add_picture(f"{SS}/03b_right_corr.png", sx + Inches(0.15), sy + Inches(3.2), Inches(2.7), Inches(2.2))

    if extra:
        add_text(slide3, sx + Inches(0.3), sy + Inches(3.2), Inches(2.4), Inches(0.8),
                 extra, size=10, color=ACCENT)

print("Slide 3: Steps done")


# ═══════════════════════════════════════════
# Slide 4: Sensor data
# ═══════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, NAVY)
add_shape(slide4, 0, 0, SW, Inches(0.08), fill_color=TEAL)

add_text(slide4, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
         "センサーデータの見方", size=28, color=WHITE, bold=True)
add_text(slide4, Inches(0.5), Inches(0.65), Inches(12), Inches(0.4),
         "4つのセンサー項目で環境の快適さを数値化します", size=14, color=LIGHT_GRAY)

# Legend screenshot
slide4.shapes.add_picture(f"{SS}/06_legend.png", Inches(0.5), Inches(1.3), Inches(2.5), Inches(3.3))

# Sensor cards
sensors = [
    ("温度", "17〜28℃", "気温を測定。快適範囲は17〜28℃。\n緯度が高い（北の）施設ほど低温、\n標高が高い施設ほど低温の傾向。", "#3AABA8"),
    ("湿度", "30〜70%", "空気中の水分量。快適範囲は30〜70%。\n夏は高湿、冬は低湿の傾向。\n沿岸部は内陸部より高い。", "#2E8EC4"),
    ("照度", "300〜750lx", "明るさを測定。快適範囲は300〜750lx。\n昼が最も高く、夜は最低。\n季節では夏が最も明るい。", "#E88A0A"),
    ("CO2", "350〜1000ppm", "二酸化炭素濃度。快適範囲は350〜1000ppm。\n夜間は高め、昼間は光合成で低下。\n標高が高い施設は低い傾向。", "#5dc573"),
]

for i, (name, range_val, desc, color) in enumerate(sensors):
    sx = Inches(3.5) + Inches(2.45) * i
    sy = Inches(1.3)
    card = add_rounded_rect(slide4, sx, sy, Inches(2.3), Inches(3.3), fill_color=RGBColor(0x16, 0x24, 0x4A))

    # Color indicator
    indicator = add_shape(slide4, sx, sy, Inches(2.3), Inches(0.06),
                         fill_color=RGBColor(int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)))

    add_text(slide4, sx + Inches(0.2), sy + Inches(0.2), Inches(2), Inches(0.5),
             name, size=20, color=WHITE, bold=True)
    add_text(slide4, sx + Inches(0.2), sy + Inches(0.65), Inches(2), Inches(0.4),
             f"快適範囲: {range_val}", size=11, color=ACCENT)
    add_text(slide4, sx + Inches(0.2), sy + Inches(1.1), Inches(1.9), Inches(2),
             desc, size=11, color=LIGHT_GRAY)

# Scoring explanation
sy2 = Inches(4.8)
add_text(slide4, Inches(0.5), sy2, Inches(12), Inches(0.5),
         "快適スコアの算出方法", size=18, color=WHITE, bold=True)
score_desc = ("各センサー値が快適範囲の中央にどれだけ近いかを0〜100で数値化し、"
              "4項目の平均を「快適スコア」として算出します。")
add_text(slide4, Inches(0.5), sy2 + Inches(0.5), Inches(8), Inches(0.5),
         score_desc, size=12, color=LIGHT_GRAY)

# Score color reference
scores = [
    ("75〜100", "快適", TEAL), ("55〜74", "良好", RGBColor(0x5D, 0xC5, 0x73)),
    ("35〜54", "普通", RGBColor(0xF0, 0xC0, 0x40)), ("0〜34", "不快", RGBColor(0xE5, 0x5A, 0x4E)),
]
for i, (rng, label, col) in enumerate(scores):
    bx = Inches(0.5) + Inches(2.5) * i
    by = sy2 + Inches(1.1)
    dot = slide4.shapes.add_shape(MSO_SHAPE.OVAL, bx, by + Inches(0.05), Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = col
    dot.line.fill.background()
    add_text(slide4, bx + Inches(0.35), by, Inches(2), Inches(0.35),
             f"{rng}  {label}", size=13, color=WHITE)

print("Slide 4: Sensor data done")


# ═══════════════════════════════════════════
# Slide 5: Popup details
# ═══════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, NAVY)
add_shape(slide5, 0, 0, SW, Inches(0.08), fill_color=TEAL)

add_text(slide5, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
         "施設の詳細を見る", size=28, color=WHITE, bold=True)
add_text(slide5, Inches(0.5), Inches(0.65), Inches(12), Inches(0.4),
         "マーカーをクリックして施設ごとの環境データを確認", size=14, color=LIGHT_GRAY)

# Popup screenshot (large)
slide5.shapes.add_picture(f"{SS}/05_popup.png", Inches(0.8), Inches(1.5), Inches(5.0), Inches(4.0))

# Explanation on the right
explanations = [
    ("施設名・種別バッジ", "施設の正式名称と「交流の家」「自然の家」の\n種別がバッジで表示されます。"),
    ("位置・標高情報", "緯度（N35.29など）と標高（500mなど）が\n表示されます。データ分析の参考に。"),
    ("快適スコア", "4つのセンサー値から算出された\n総合的な快適度が色付きで表示されます。"),
    ("センサーバー", "温度・湿度・照度・CO2の各値が\nバーグラフで視覚的に表示されます。\n色はスコアに連動しています。"),
]

ey = Inches(1.5)
for title, desc in explanations:
    # Arrow/bullet
    add_text(slide5, Inches(6.2), ey, Inches(6), Inches(0.4),
             title, size=16, color=ACCENT, bold=True)
    add_text(slide5, Inches(6.2), ey + Inches(0.4), Inches(6), Inches(0.8),
             desc, size=12, color=LIGHT_GRAY)
    ey += Inches(1.3)

# Sidebar screenshot at bottom
slide5.shapes.add_picture(f"{SS}/02_sidebar.png", Inches(7.5), Inches(1.5), Inches(3.0), Inches(5.5))

# Note about sidebar
add_text(slide5, Inches(10.8), Inches(2.0), Inches(2.3), Inches(3),
         "左パネルの施設一覧からも\n選択できます。\n\nクリックすると地図が\nその施設に移動し、\nポップアップが表示されます。",
         size=11, color=LIGHT_GRAY)

print("Slide 5: Popup done")


# ═══════════════════════════════════════════
# Slide 6: Data analysis (right panel)
# ═══════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, NAVY)
add_shape(slide6, 0, 0, SW, Inches(0.08), fill_color=TEAL)

add_text(slide6, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
         "データの分析方法", size=28, color=WHITE, bold=True)
add_text(slide6, Inches(0.5), Inches(0.65), Inches(12), Inches(0.4),
         "右パネルの6つのタブで多角的にデータを分析できます", size=14, color=LIGHT_GRAY)

# Three panel screenshots
panels = [
    (f"{SS}/03_right_season.png", "季節タブ", "現在の季節・時間帯の詳細情報、\n全施設の平均値、\nデータサイエンスのヒントを表示"),
    (f"{SS}/03c_right_dist.png", "分布タブ", "ヒストグラムと箱ひげ図で\nデータの分布を可視化。\n基本統計量（平均/中央値/標準偏差）も表示"),
    (f"{SS}/03b_right_corr.png", "相関タブ", "散布図で2変数間の関係を分析。\n相関係数と回帰直線を自動計算。\n「緯度×温度」で地理的傾向を発見"),
]

for i, (img_path, title, desc) in enumerate(panels):
    px = Inches(0.3) + Inches(4.3) * i
    py = Inches(1.2)
    slide6.shapes.add_picture(img_path, px, py, Inches(3.5), Inches(3.8))
    add_text(slide6, px, py + Inches(3.9), Inches(3.5), Inches(0.4),
             title, size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide6, px, py + Inches(4.3), Inches(3.5), Inches(1.2),
             desc, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Additional tabs explanation
ay = Inches(6.2)
add_shape(slide6, Inches(0.3), ay, Inches(12.7), Inches(0.04), fill_color=BLUE2)

other_tabs = [
    ("時系列タブ", "施設ごとの一日の変化（朝→昼→夕→夜）をグラフで比較"),
    ("異常値タブ", "平均±2σ（標準偏差）から外れた値を自動検出"),
    ("CSVタブ", "全データをCSVファイルとしてダウンロード（全16パターン対応）"),
]
for i, (tab, desc) in enumerate(other_tabs):
    tx = Inches(0.5) + Inches(4.3) * i
    add_text(slide6, tx, ay + Inches(0.15), Inches(1.5), Inches(0.35),
             tab, size=13, color=ACCENT, bold=True)
    add_text(slide6, tx + Inches(1.6), ay + Inches(0.17), Inches(2.5), Inches(0.35),
             desc, size=11, color=LIGHT_GRAY)

print("Slide 6: Analysis done")


# ═══════════════════════════════════════════
# Slide 7: Access URL & tips
# ═══════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, NAVY)
add_shape(slide7, 0, 0, SW, Inches(0.08), fill_color=TEAL)

add_text(slide7, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
         "アクセス方法・授業での活用", size=28, color=WHITE, bold=True)

# URL box
url_box = add_rounded_rect(slide7, Inches(1), Inches(1.2), Inches(11.3), Inches(1.4), fill_color=RGBColor(0x16, 0x24, 0x4A))
add_text(slide7, Inches(1.3), Inches(1.3), Inches(10), Inches(0.4),
         "アクセスURL", size=14, color=ACCENT, bold=True)
add_text(slide7, Inches(1.3), Inches(1.7), Inches(10), Inches(0.6),
         "https://shiojima-test.github.io/schoomy-gis/niye_comfort_v1.html",
         size=20, color=WHITE, bold=True)

# Requirements
add_text(slide7, Inches(0.5), Inches(3.0), Inches(5.5), Inches(0.5),
         "動作環境", size=18, color=WHITE, bold=True)
reqs = [
    "ブラウザ: Chrome / Edge 推奨",
    "インターネット接続: 必要（地図タイル読み込み）",
    "インストール: 不要（Webアプリ）",
    "シリアル接続: Chrome/Edge の Web Serial API 対応",
    "（実機センサー接続はオプション）",
]
ry = Inches(3.5)
for req in reqs:
    add_text(slide7, Inches(0.8), ry, Inches(5), Inches(0.35),
             f"  {req}", size=12, color=LIGHT_GRAY)
    ry += Inches(0.35)

# Activity ideas
add_text(slide7, Inches(7), Inches(3.0), Inches(5.5), Inches(0.5),
         "授業での活用アイデア", size=18, color=WHITE, bold=True)
ideas = [
    ("探究1", "「なぜ北の施設ほど気温が低いのか？」\n緯度と温度の相関を散布図で確認"),
    ("探究2", "「標高が環境に与える影響は？」\n標高と各センサー値の関係を分析"),
    ("探究3", "「季節による環境変化のパターンは？」\n4季節×4時間帯のデータを比較"),
    ("探究4", "「異常値のある施設の特徴は？」\n異常値タブから外れ値の原因を推測"),
]
iy = Inches(3.5)
for num, desc in ideas:
    badge = add_rounded_rect(slide7, Inches(7), iy, Inches(1), Inches(0.35), fill_color=TEAL)
    tf = badge.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Meiryo"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(slide7, Inches(8.2), iy - Inches(0.05), Inches(4.5), Inches(0.9),
             desc, size=11, color=LIGHT_GRAY)
    iy += Inches(0.9)

# Bottom
add_shape(slide7, 0, SH - Inches(0.6), SW, Inches(0.6), fill_color=BLUE2)
add_text(slide7, Inches(1), SH - Inches(0.55), Inches(11), Inches(0.5),
         "国立青少年教育施設 環境マッピングマップ  |  高校情報II データサイエンス授業用",
         size=11, color=WHITE, align=PP_ALIGN.CENTER)

print("Slide 7: Access done")


# Save
prs.save(OUT)
print(f"\nPowerPoint saved: {OUT}")
